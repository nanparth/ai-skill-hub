"""PPTX adapter: parse slide titles and bullet bodies into NormalizedDoc.

Heavy import (python-pptx) is LAZY: imported inside parse(), never at module top.
Note: slide images are not parsed; only text frames are extracted.
"""
from __future__ import annotations

from . import NormalizedDoc, DocBlock, limit_value, mark_truncated


def parse(src: str, **opts) -> NormalizedDoc:
    from pptx import Presentation  # lazy heavy import

    doc = NormalizedDoc(source_format="pptx")
    prs = Presentation(src)
    max_slides = limit_value(opts, "max_pptx_slides")
    max_text_shapes = limit_value(opts, "max_pptx_text_shapes")

    idx = 0
    text_shapes_seen = 0
    for slide_no, slide in enumerate(prs.slides, start=1):
        if max_slides and slide_no > max_slides:
            mark_truncated(doc, "PPTX_SLIDE_LIMIT_REACHED")
            break
        anchor = f"slide{slide_no}"

        # Title placeholder text -> heading block.
        title_text = None
        title_shape = getattr(slide.shapes, "title", None)
        title_id = None
        if title_shape is not None and title_shape.has_text_frame:
            title_text = title_shape.text_frame.text.strip() or None
            title_id = getattr(title_shape, "shape_id", None)

        slide_heading_path = [title_text] if title_text else []

        if title_text:
            doc.blocks.append(
                DocBlock(
                    text=title_text,
                    block_type="heading",
                    idx=idx,
                    anchor=anchor,
                    parent_heading=None,
                    heading_path=slide_heading_path,
                )
            )
            idx += 1

        # Body shapes -> list_item blocks (each paragraph is a bullet).
        for shape in slide.shapes:
            if title_id is not None and getattr(shape, "shape_id", None) == title_id:
                continue
            if not getattr(shape, "has_text_frame", False):
                continue  # slide images and non-text shapes skipped
            if max_text_shapes and text_shapes_seen >= max_text_shapes:
                mark_truncated(doc, "PPTX_TEXT_SHAPE_LIMIT_REACHED")
                break
            text_shapes_seen += 1
            for para in shape.text_frame.paragraphs:
                bullet = "".join(run.text for run in para.runs).strip()
                if not bullet:
                    bullet = (para.text or "").strip()
                if not bullet:
                    continue
                doc.blocks.append(
                    DocBlock(
                        text=bullet,
                        block_type="list_item",
                        idx=idx,
                        anchor=anchor,
                        parent_heading=title_text,
                        heading_path=list(slide_heading_path),
                    )
                )
                idx += 1

    return doc
